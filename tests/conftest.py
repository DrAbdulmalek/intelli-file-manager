"""إعدادات مشتركة للاختبارات -fixtures وmocks"""
import sys
import os
import json
import tempfile
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest

# تأكد من أن مسار المشروع الجذر في sys.path
PROJECT_ROOT = Path(__file__).resolve().parent.parent
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from src.core.config import Config, CATEGORIES, EXTENSION_MAP


# ─── Fixture: مجلد مؤقت مع ملفات عيّنة ────────────────────────────────────
@pytest.fixture
def tmp_dir(tmp_path):
    """ينشئ مجلداً مؤقتاً يحتوي ملفات متنوعة للاختبار

    يحتوي على:
      - مستندات: report.pdf, notes.txt, data.csv
      - صور: photo.jpg, image.png
      - برمجة: script.py, app.js
      - أرشيفات: backup.zip
      - ملفات بامتدادات غير معروفة: file.xyz123
    """
    # ملفات المستندات
    (tmp_path / "report.pdf").write_bytes(b"%PDF-1.4 fake pdf content")
    (tmp_path / "notes.txt").write_text("هذه ملاحظات تجريبية", encoding="utf-8")
    (tmp_path / "data.csv").write_text("اسم,عمر\nأحمد,25\nسارة,30\n", encoding="utf-8")

    # صور
    (tmp_path / "photo.jpg").write_bytes(b"\xff\xd8\xff\xe0 JFIF fake jpeg")
    (tmp_path / "image.png").write_bytes(b"\x89PNG\r\n\x1a\n fake png")

    # برمجة
    (tmp_path / "script.py").write_text("print('مرحباً بالعالم')\n", encoding="utf-8")
    (tmp_path / "app.js").write_text("console.log('hello');\n", encoding="utf-8")

    # أرشيفات
    (tmp_path / "backup.zip").write_bytes(b"PK\x03\x04 fake zip")

    # ملف بامتداد غير معروف
    (tmp_path / "file.xyz123").write_bytes(b"unknown content")

    # مجلد فرعي مع ملفات
    subdir = tmp_path / "مجلد_فرعي"
    subdir.mkdir()
    (subdir / "document.docx").write_bytes(b"PK fake docx")
    (subdir / "song.mp3").write_bytes(b"ID3 fake mp3")

    return tmp_path


# ─── Fixture: كائن Config بإعدادات اختبار ────────────────────────────────
@pytest.fixture
def sample_config(tmp_path):
    """ينشئ كائن Config بإعدادات اختبار

    يستخدم مسار قاعدة بيانات مؤقت ويضيف تصنيفاً مخصصاً للتجربة.
    """
    config = Config(
        database_path=str(tmp_path / "intellifile_test"),
        ai_model="gemma3",
        ollama_url="http://localhost:11434",
        language="ar",
        dark_mode=True,
        auto_classify=True,
        duplicate_detection=True,
        file_protection=True,
        voice_enabled=False,
    )
    # إضافة تصنيف مخصص للاختبار
    config.add_custom_category("اختبارات", [".xyz123", ".tdata"])
    return config


# ─── Fixture: ملفات عيّنة إضافية ─────────────────────────────────────────
@pytest.fixture
def sample_files(tmp_path):
    """ينشئ ملفات عيّنة متنوعة في المجلد المؤقت ويُرجع قاموساً بالمسارات

    يُرجع قاموساً يحتوي على:
      - pdf_file: مسار ملف PDF
      - txt_file: مسار ملف نصي
      - py_file: مسار ملف Python
      - jpg_file: مسار صورة JPEG
      - zip_file: مسار أرشيف ZIP
      - unknown_file: ملف بامتداد غير معروف
      - empty_file: ملف فارغ
      - nested_dir: مجلد فرعي
      - nested_file: ملف داخل المجلد الفرعي
    """
    pdf_file = tmp_path / "document.pdf"
    pdf_file.write_bytes(b"%PDF-1.4 test content")

    txt_file = tmp_path / "readme.txt"
    txt_file.write_text("محتوى نصي للتجربة\n", encoding="utf-8")

    py_file = tmp_path / "main.py"
    py_file.write_text("def hello():\n    print('مرحباً')\n", encoding="utf-8")

    jpg_file = tmp_path / "avatar.jpg"
    jpg_file.write_bytes(b"\xff\xd8\xff\xe0 fake jpeg data")

    zip_file = tmp_path / "archive.zip"
    zip_file.write_bytes(b"PK\x03\x04 fake archive")

    unknown_file = tmp_path / "data.unknown_ext"
    unknown_file.write_bytes(b"unknown data content")

    empty_file = tmp_path / "empty.txt"
    empty_file.write_bytes(b"")

    nested_dir = tmp_path / "subfolder"
    nested_dir.mkdir()
    nested_file = nested_dir / "nested.py"
    nested_file.write_text("# ملف متداخل\npass\n", encoding="utf-8")

    return {
        "pdf_file": str(pdf_file),
        "txt_file": str(txt_file),
        "py_file": str(py_file),
        "jpg_file": str(jpg_file),
        "zip_file": str(zip_file),
        "unknown_file": str(unknown_file),
        "empty_file": str(empty_file),
        "nested_dir": str(nested_dir),
        "nested_file": str(nested_file),
        "tmp_path": str(tmp_path),
    }


# ─── Fixture: ملفات حقيقية قابلة للاستخراج (txt/pdf/docx/xlsx/pptx) ───────
# تُستخدم لاختبارات FileInventory التكاملية. تُبنى الملفات فعليًا باستخدام
# المكتبات (pdfplumber, python-docx, openpyxl, python-pptx) بدلًا من fake bytes.
@pytest.fixture
def real_doc_dir(tmp_path):
    """ينشئ مجلداً يحتوي ملفات حقيقية من الأنواع الخمسة المدعومة

    المحتوى عربي + إنجليزي لاختبار الترميز. كل ملف يحتوي نصًا قابلًا للاستخراج.
    """
    # ─── txt: نص مباشر ───
    (tmp_path / "notes.txt").write_text(
        "هذه ملاحظات تجريبية\nThis is a test note\nخط ثالث",
        encoding="utf-8",
    )

    # ─── md: نص مباشر ───
    (tmp_path / "readme.md").write_text(
        "# عنوان\n\nفقرة وصفية بالعربية.",
        encoding="utf-8",
    )

    # ─── pdf: يُبنى عبر مكتبة reportlab إن وُجدت، وإلا يُكتب نص فارغ ───
    _build_real_pdf(tmp_path / "report.pdf")

    # ─── docx: يُبنى عبر python-docx ───
    _build_real_docx(tmp_path / "document.docx")

    # ─── xlsx: يُبنى عبر openpyxl ───
    _build_real_xlsx(tmp_path / "spreadsheet.xlsx")

    # ─── pptx: يُبنى عبر python-pptx ───
    _build_real_pptx(tmp_path / "presentation.pptx")

    # ─── ملف مكرر (نسخة من notes.txt) ───
    (tmp_path / "notes_copy.txt").write_text(
        "هذه ملاحظات تجريبية\nThis is a test note\nخط ثالث",
        encoding="utf-8",
    )

    # ─── ملف مخفي (يجب تخطيه) ───
    (tmp_path / ".hidden_secret").write_text("secret", encoding="utf-8")

    # ─── مجلد فرعي مع ملف ───
    sub = tmp_path / "subfolder"
    sub.mkdir()
    (sub / "nested.txt").write_text("نص متداخل", encoding="utf-8")

    # ─── ملف غير مدعوم للاستخراج (jpg ببيانات وهمية) ───
    (tmp_path / "image.jpg").write_bytes(b"\xff\xd8\xff\xe0 fake")

    return tmp_path


def _build_real_pdf(path):
    """يبني PDF حقيقي بصفحة واحدة تحتوي نصًا عربيًا وإنجليزيًا"""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
        c = canvas.Canvas(str(path), pagesize=A4)
        c.drawString(100, 750, "Hello from PDF")
        c.drawString(100, 730, "Test content for extraction")
        c.drawString(100, 710, "Arabic: ahlan wa sahlan")
        c.showPage()
        c.save()
    except ImportError:
        # fallback: اكتب PDF بسيط يدويًا (نص قابل للاستخراج)
        # هذا PDF بسيط جدًا بصفحة واحدة — pdfplumber قد يستخرج نصًا منه
        path.write_bytes(_minimal_pdf_bytes())


def _minimal_pdf_bytes() -> bytes:
    """PDF بسيط بصفحة واحدة يحتوي نصًا يمكن استخراجه"""
    content = b"BT /F1 12 Tf 100 700 Td (Hello from PDF) Tj ET"
    return (
        b"%PDF-1.4\n"
        b"1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj\n"
        b"2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj\n"
        b"3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >> endobj\n"
        b"4 0 obj << /Length " + str(len(content)).encode() + b" >> stream\n" +
        content + b"\nendstream endobj\n"
        b"5 0 obj << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> endobj\n"
        b"xref\n0 6\n0000000000 65535 f \n"
        b"trailer << /Size 6 /Root 1 0 R >>\nstartxref\n0\n%%EOF\n"
    )


def _build_real_docx(path):
    """يبني DOCX حقيقي بفقرات وجدول"""
    from docx import Document
    doc = Document()
    doc.add_paragraph("Hello from DOCX")
    doc.add_paragraph("فقرة عربية تجريبية")
    doc.add_paragraph("Test content for extraction")
    # جدول بسيط
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Age"
    table.cell(1, 0).text = "Ahmad"
    table.cell(1, 1).text = "30"
    doc.save(str(path))


def _build_real_xlsx(path):
    """يبني XLSX حقيقي بورقتين وبيانات"""
    import openpyxl
    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Sheet1"
    ws1.append(["Name", "Age", "City"])
    ws1.append(["Ahmad", 30, "Damascus"])
    ws1.append(["Sara", 28, "Cairo"])
    # ورقة ثانية
    ws2 = wb.create_sheet("Sheet2")
    ws2.append(["Product", "Price"])
    ws2.append(["Book", 15.5])
    wb.save(str(path))


def _build_real_pptx(path):
    """يبني PPTX حقيقي بشريحتين تحتويان نصًا وجدولًا"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    # شريحة 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "First Slide"
    slide1.placeholders[1].text = "Hello from PPTX\nArabic content"
    # شريحة 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    left = top = width = height = Inches(2)
    txBox = slide2.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Second slide content"
    # جدول في الشريحة 2
    table_shape = slide2.shapes.add_table(
        rows=2, cols=2, left=Inches(1), top=Inches(3),
        width=Inches(4), height=Inches(1)
    )
    table_shape.table.cell(0, 0).text = "Col1"
    table_shape.table.cell(0, 1).text = "Col2"
    table_shape.table.cell(1, 0).text = "Val1"
    table_shape.table.cell(1, 1).text = "Val2"
    prs.save(str(path))


# ─── Fixture: ملفات وسائط حقيقية (jpg/png/mp3/mp4) لاختبارات PR-03 ───────
# تُبنى فعليًا عبر PIL (للصور) وffmpeg (للصوت/الفيديو).
@pytest.fixture
def real_media_dir(tmp_path):
    """ينشئ مجلداً يحتوي ملفات وسائط حقيقية لاختبار الميتاداتا الموسّعة

    المحتوى:
      - sample.jpg: JPEG 80x60 مع EXIF (Make, Model, DateTime)
      - sample.png: PNG 50x50
      - sample.mp3: ملف صوتي 1.5s صمت (يُولَّد عبر ffmpeg)
      - sample.mp4: ملف فيديو 1s لون ثابت (يُولَّد عبر ffmpeg)
      - corrupt.jpg: ملف بامتداد jpg لكن بمحتوى عشوائي (لاختبار التسامح)
    """
    _build_real_jpeg_with_exif(tmp_path / "sample.jpg")
    _build_real_png(tmp_path / "sample.png")
    _build_real_mp3(tmp_path / "sample.mp3")
    _build_real_mp4(tmp_path / "sample.mp4")
    # ملف تالف لاختبار التسامح مع الأخطاء
    (tmp_path / "corrupt.jpg").write_bytes(b"\xff\xd8\xff\xe0 not really a jpeg")
    return tmp_path


def _build_real_jpeg_with_exif(path):
    """يبني JPEG حقيقي 80x60 مع وسوم EXIF (Make, Model, DateTime)"""
    from PIL import Image
    img = Image.new("RGB", (80, 60), color=(123, 200, 50))
    exif = img.getexif() if hasattr(img, "getexif") else None
    if exif is not None:
        # وسوم EXIF الشائعة
        # 0x010F = Make, 0x0110 = Model, 0x0132 = DateTime
        exif[0x010F] = "TestCamera"
        exif[0x0110] = "IFM-1000"
        exif[0x0132] = "2026:07:24 10:30:00"
        exif[0x013B] = "Artist Name"  # Artist
    img.save(str(path), "JPEG", exif=exif)


def _build_real_png(path):
    """يبني PNG حقيقي 50x50"""
    from PIL import Image
    img = Image.new("RGBA", (50, 50), color=(10, 20, 30, 255))
    img.save(str(path), "PNG")


def _build_real_mp3(path):
    """يبني MP3 حقيقي 1.5s صمت عبر ffmpeg (fallback: يكتب بايتات وهمية)"""
    import subprocess
    import shutil
    ffmpeg = shutil.which("ffmpeg")
    if not ffmpeg:
        # fallback: لا يوجد ffmpeg، نكتب رأس MP3 فقط
        path.write_bytes(b"ID3\x03\x00\x00\x00\x00\x00\x00\x00" + b"\xff\xfb" * 100)
        return
    try:
        subprocess.run(
            [ffmpeg, "-y", "-f", "lavfi", "-i", "anullsrc=r=44100:cl=mono",
             "-t", "1.5", "-q:a", "9", "-metadata",
             "title=Test Audio", "-metadata", "artist=IFM Test",
             str(path)],
            capture_output=True, timeout=15, check=False,
        )
    except Exception:
        path.write_bytes(b"ID3\x03\x00\x00\x00\x00\x00\x00\x00" + b"\xff\xfb" * 100)


def _build_real_mp4(path):
    """يبني MP4 حقيقي 1s لون ثابت عبر ffmpeg (fallback: يكتب رأس MP4)"""
    import subprocess
    import shutil
    ffmpeg = shutil.which("ffmpeg")
    if not ffmpeg:
        # fallback: رأس MP4 بسيط
        path.write_bytes(b"\x00\x00\x00\x18 ftypmp42" + b"\x00" * 50)
        return
    try:
        subprocess.run(
            [ffmpeg, "-y", "-f", "lavfi", "-i",
             "color=c=red:s=160x120:d=1:r=15",
             "-c:v", "libx264", "-pix_fmt", "yuv420p",
             "-movflags", "+faststart",
             str(path)],
            capture_output=True, timeout=20, check=False,
        )
    except Exception:
        path.write_bytes(b"\x00\x00\x00\x18 ftypmp42" + b"\x00" * 50)


# ─── Mock: ollama ──────────────────────────────────────────────────────────
@pytest.fixture
def mock_ollama():
    """يُحاكي مكتبة ollama بدون اتصال حقيقي بالخدمة

    يوفر كائن chat الاستجابات التالية:
      - وصف لمحتوى ملف
    """
    mock = MagicMock()
    mock.chat.return_value = "هذا ملف مستند نصي يحتوي على معلومات مهمة"
    mock.list.return_value = {"models": [{"name": "gemma3"}]}
    mock.show.return_value = {
        "details": {"parameter_size": "2B", "quantization_level": "Q4_0"}
    }
    return mock


# ─── Mock: chromadb ────────────────────────────────────────────────────────
@pytest.fixture
def mock_chromadb():
    """يُحاكي مكتبة chromadb بدون اتصال حقيقي بقاعدة البيانات

    يوفر:
      - PersistentClient: عميل وهمي
      - collection: مجموعة وهمية للبيانات
    """
    mock_client = MagicMock()
    mock_collection = MagicMock()
    mock_collection.add.return_value = ["id1", "id2"]
    mock_collection.query.return_value = {
        "ids": [["id1"]],
        "documents": [["نص تجريبي"]],
        "metadatas": [[{"source": "test.txt"}]],
        "distances": [[0.1]],
    }
    mock_client.get_or_create_collection.return_value = mock_collection

    with patch("chromadb.PersistentClient", return_value=mock_client):
        yield {
            "client": mock_client,
            "collection": mock_collection,
        }


# ─── Mock: magika ──────────────────────────────────────────────────────────
@pytest.fixture
def mock_magika():
    """يُحاكي مكتبة magika للتصنيف بالمحتوى

    ينشئ كائن وهمي يدعم أسلوب identify_path الذي يُرجع
    كائناً يحتوي على خاصية output بنوع الملف.
    """
    mock_result = MagicMock()
    mock_result.output = "pdf"

    mock_magika_instance = MagicMock()
    mock_magika_instance.identify_path.return_value = mock_result

    with patch.dict("sys.modules", {"magika": MagicMock(Magika=lambda: mock_magika_instance)}):
        yield mock_magika_instance


# ─── Fixture: منع وحدات خارجية عن طريق التصحيح ──────────────────────────
@pytest.fixture(autouse=True)
def patch_external_modules():
    """يمنع تلقائياً تحميل الوحدات الخارجية التي قد لا تكون مثبتة

    يشمل: ollama, chromadb, magika, speech_recognition, pyttsx3
    """
    modules_to_mock = [
        "ollama",
        "chromadb",
        "magika",
        "speech_recognition",
        "pyttsx3",
    ]
    mocked = {}
    for mod in modules_to_mock:
        if mod not in sys.modules:
            mocked[mod] = MagicMock()
            sys.modules[mod] = mocked[mod]

    yield

    # تنظيف: إزالة فقط ما أضفناه
    for mod in modules_to_mock:
        if mod in mocked and mod in sys.modules:
            del sys.modules[mod]
