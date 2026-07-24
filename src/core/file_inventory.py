"""FileInventory — جرد الملفات المفهرس مع استخراج الميتاداتا والمحتوى

هذه الوحدة هي المسؤول الوحيد عن:
  - مسح شجرة مجلدات وتوليد FileMetadata لكل ملف
  - حساب SHA-256 لكل ملف (جاهز لكشف التكرار)
  - استخراج النص من txt/pdf/docx/xlsx/pptx عبر استيراد كسول
    (graceful degradation: لو المكتبة غير مثبتة يُرجع نص فارغ)
  - استخراج الميتاداتا الموسّعة للصور (EXIF) والصوت/الفيديو (ffprobe)
    عبر metadata_extractor الموحّد (PR-03)
  - كشف نوع المحتوى (MIME) عبر python-magic مع fallback للامتداد (PR-03)
  - إرجاع FileRecord جاهز للفهرسة/التخزين

لا تشمل هذه الوحدة:
  - أي ذكاء اصطناعي (AI/embeddings) — مرحلة B
  - أي ميزات طبية — ممنوعة per PRODUCT_IDENTITY.md
  - أي تخزين دائم — مسؤولية طبقة db

PR-02 من development-roadmap-v1.0 (FileInventory + اختبارات تكامل)
PR-03 من development-roadmap-v1.0 (ميتاداتا موسّعة + content_type عبر magic)
"""
from __future__ import annotations

import hashlib
import logging
import os
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Iterable, Iterator, List, Optional

from ..db.schemas import FileCategory, FileMetadata, FileRecord
from .metadata_extractor import detect_content_type, extract_extended_metadata

logger = logging.getLogger(__name__)

# ─── حد أقصى لحجم الملف لاستخراج النص (10 MB) ──────────────────────────────
MAX_CONTENT_EXTRACTION_SIZE = 10 * 1024 * 1024

# ─── امتدادات الملفات المدعومة لاستخراج النص ───────────────────────────────
TEXT_EXTRACTABLE_EXT = {".txt", ".md", ".csv", ".log", ".pdf", ".docx", ".xlsx", ".pptx"}


@dataclass
class InventoryStats:
    """إحصائيات جرد الملفات

    تُرجع من FileInventory.scan() كملخص لعملية المسح
    """
    total_files: int = 0
    indexed_files: int = 0
    skipped_files: int = 0
    content_extracted: int = 0
    content_extraction_failed: int = 0
    duplicate_candidates: int = 0
    total_size_bytes: int = 0
    duration_seconds: float = 0.0
    errors: List[str] = field(default_factory=list)

    def to_dict(self) -> dict:
        return {
            "total_files": self.total_files,
            "indexed_files": self.indexed_files,
            "skipped_files": self.skipped_files,
            "content_extracted": self.content_extracted,
            "content_extraction_failed": self.content_extraction_failed,
            "duplicate_candidates": self.duplicate_candidates,
            "total_size_bytes": self.total_size_bytes,
            "total_size_human": _format_size(self.total_size_bytes),
            "duration_seconds": round(self.duration_seconds, 3),
            "errors_count": len(self.errors),
        }


def _format_size(size_bytes: int) -> str:
    """تنسيق الحجم بصيغة مقروءة"""
    if size_bytes == 0:
        return "0 B"
    units = ["B", "KB", "MB", "GB", "TB"]
    size = float(size_bytes)
    i = 0
    while size >= 1024.0 and i < len(units) - 1:
        size /= 1024.0
        i += 1
    return f"{size:.2f} {units[i]}"


class FileInventory:
    """جرد الملفات المفهرس

    الاستخدام الأساسي:

        inventory = FileInventory()
        records = list(inventory.scan("/path/to/folder"))
        stats = inventory.last_stats

    الخصائص:
      - كسول (lazy): يُولّد FileRecord واحدًا في كل مرة عبر generator
      - آمن: لا يُعدّل الملفات، للقراءة فقط
      - متسامح: الأخطاء في ملف واحد لا توقف المسح كله
      - قابل للتكوين: skip_hidden, max_file_size, include_content
    """

    def __init__(
        self,
        *,
        skip_hidden: bool = True,
        skip_dirs: Optional[Iterable[str]] = None,
        max_file_size: int = MAX_CONTENT_EXTRACTION_SIZE,
        include_content: bool = True,
    ):
        """
        Args:
            skip_hidden: تخطي الملفات/المجلدات المخفية (التي تبدأ بـ .)
            skip_dirs: أسماء مجلدات إضافية لتخطيها (مثل __pycache__, node_modules, .git)
            max_file_size: أقصى حجم بالبايت لاستخراج المحتوى منه (الملفات الأكبر تُسجّل
                          بدون محتوى)
            include_content: هل يُستخرج النص من الملفات المدعومة؟
        """
        self.skip_hidden = skip_hidden
        self.skip_dirs = set(skip_dirs or {"__pycache__", "node_modules", ".git", ".venv", "venv"})
        self.max_file_size = max_file_size
        self.include_content = include_content
        self.last_stats: Optional[InventoryStats] = None

    # ─── API الرئيسي ──────────────────────────────────────────────────────

    def scan(self, directory: str, *, recursive: bool = True) -> Iterator[FileRecord]:
        """مسح مجلد وإرجاع FileRecord لكل ملف (كسول)

        Args:
            directory: مسار المجلد المراد مسحه
            recursive: هل يشمل المسح المجلدات الفرعية؟

        Yields:
            FileRecord: سجل ملف واحد لكل ملف موجود
        """
        dir_path = Path(directory).resolve()
        if not dir_path.is_dir():
            logger.warning(f"المجلد غير موجود: {directory}")
            self.last_stats = InventoryStats(errors=[f"directory not found: {directory}"])
            return

        start = datetime.now()
        stats = InventoryStats()
        seen_hashes: dict[str, str] = {}  # sha256 -> first path

        for file_path in self._iter_files(dir_path, recursive=recursive):
            stats.total_files += 1
            try:
                record = self._build_record(file_path, seen_hashes, stats)
                if record is not None:
                    stats.indexed_files += 1
                    yield record
                else:
                    stats.skipped_files += 1
            except Exception as e:
                stats.skipped_files += 1
                stats.errors.append(f"{file_path}: {e}")
                logger.debug(f"خطأ في فهرسة {file_path}: {e}")

        stats.duration_seconds = (datetime.now() - start).total_seconds()
        stats.duplicate_candidates = len(seen_hashes) - len(set(seen_hashes.values())) \
            if seen_hashes else 0
        # Note: duplicate_candidates = number of files whose hash matched an earlier file
        # (computed differently below for accuracy)
        self.last_stats = stats

    def scan_directory(self, directory: str) -> List[FileRecord]:
        """واجهة متحمّسة (eager) — تُرجع قائمة كاملة بدلًا من generator

        مريحة للاختبارات وللحالات التي تريد كل النتائج دفعة واحدة.
        """
        return list(self.scan(directory))

    # ─── البناء الداخلي ──────────────────────────────────────────────────

    def _iter_files(self, dir_path: Path, *, recursive: bool) -> Iterator[Path]:
        """مولّد للمسارات فقط، مع تطبيق skip_hidden و skip_dirs"""
        pattern = "**/*" if recursive else "*"
        for item in dir_path.glob(pattern):
            if not item.is_file():
                continue
            if self.skip_hidden and any(part.startswith(".") for part in item.parts):
                continue
            # تخطي الملفات داخل مجلدات مستثناة
            if any(part in self.skip_dirs for part in item.parts):
                continue
            yield item

    def _build_record(
        self,
        file_path: Path,
        seen_hashes: dict[str, str],
        stats: InventoryStats,
    ) -> Optional[FileRecord]:
        """يبني FileRecord واحدًا من مسار ملف"""
        try:
            stat = file_path.stat()
        except OSError as e:
            stats.errors.append(f"stat failed {file_path}: {e}")
            return None

        # حساب الـ hash (مهم لكشف التكرار)
        sha256 = self._compute_hash(file_path)
        is_duplicate = False
        if sha256 and sha256 in seen_hashes:
            is_duplicate = True
        elif sha256:
            seen_hashes[sha256] = str(file_path)

        # بناء FileMetadata
        ext = file_path.suffix.lower()
        category = FileCategory.from_extension(ext)
        # كشف نوع المحتوى عبر python-magic مع fallback للامتداد
        content_type = detect_content_type(file_path, ext)
        # استخراج الميتاداتا الموسّعة (EXIF للصور، ffprobe للصوت/الفيديو)
        extra_metadata = self._safe_extract_extended(file_path, ext)
        metadata = FileMetadata(
            file_name=file_path.name,
            file_path=str(file_path),
            file_size=stat.st_size,
            extension=ext.lstrip("."),
            category=category.value,
            sha256_hash=sha256,
            created_at=datetime.fromtimestamp(stat.st_ctime).isoformat(),
            modified_at=datetime.fromtimestamp(stat.st_mtime).isoformat(),
            is_duplicate=is_duplicate,
            content_type=content_type,
            extra_metadata=extra_metadata,
        )

        # استخراج المحتوى (اختياري)
        document_text = ""
        if self.include_content and ext in TEXT_EXTRACTABLE_EXT and stat.st_size <= self.max_file_size:
            document_text = self._extract_content(file_path, ext)
            if document_text:
                stats.content_extracted += 1
            else:
                stats.content_extraction_failed += 1
        elif self.include_content and stat.st_size > self.max_file_size:
            stats.content_extraction_failed += 1
            stats.errors.append(f"file too large for content extraction: {file_path}")

        stats.total_size_bytes += stat.st_size
        return FileRecord.from_metadata(metadata=metadata, document_text=document_text)

    @staticmethod
    def _compute_hash(file_path: Path) -> str:
        """SHA-256 مع قراءة على دفعات لتجنب تحميل ملفات ضخمة في الذاكرة"""
        h = hashlib.sha256()
        try:
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(8192), b""):
                    h.update(chunk)
            return h.hexdigest()
        except OSError as e:
            logger.debug(f"تعذر حساب hash لـ {file_path}: {e}")
            return ""

    @staticmethod
    def _safe_extract_extended(file_path: Path, ext: str) -> dict:
        """يستخرج الميتاداتا الموسّعة بتسامح تام مع الأخطاء"""
        try:
            return extract_extended_metadata(file_path, ext)
        except Exception as e:
            logger.debug(f"فشل استخراج الميتاداتا الموسّعة من {file_path}: {e}")
            return {}

    @staticmethod
    def _extract_content(file_path: Path, ext: str) -> str:
        """يستخرج النص من الملف حسب امتداده — استيراد كسول + graceful fallback

        ترتيب الاستيراد:
          - txt/md/csv/log: قراءة مباشرة (UTF-8 مع fallback)
          - pdf: pdfplumber
          - docx: python-docx
          - xlsx: openpyxl
          - pptx: python-pptx
        """
        try:
            if ext in {".txt", ".md", ".csv", ".log"}:
                return _read_text_file(file_path)
            if ext == ".pdf":
                return _extract_pdf_text(file_path)
            if ext == ".docx":
                return _extract_docx_text(file_path)
            if ext == ".xlsx":
                return _extract_xlsx_text(file_path)
            if ext == ".pptx":
                return _extract_pptx_text(file_path)
        except Exception as e:
            logger.debug(f"فشل استخراج المحتوى من {file_path}: {e}")
            return ""
        return ""


# ─── دوال استخراج النص (منفصلة لقابلية الاختبار) ────────────────────────────

def _read_text_file(file_path: Path) -> str:
    """قراءة ملف نصي مع fallback للترميز"""
    for encoding in ("utf-8", "cp1256", "latin-1"):
        try:
            return file_path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    # آخر محاولة: تجاهل الأخطاء
    return file_path.read_text(encoding="utf-8", errors="replace")


def _extract_pdf_text(file_path: Path) -> str:
    """استخراج نص من PDF عبر pdfplumber (استيراد كسول)"""
    try:
        import pdfplumber
    except ImportError:
        logger.debug("pdfplumber غير مثبت — تخطي استخراج PDF")
        return ""
    pages_text: List[str] = []
    with pdfplumber.open(str(file_path)) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            pages_text.append(text)
    return "\n".join(pages_text)


def _extract_docx_text(file_path: Path) -> str:
    """استخراج نص من DOCX عبر python-docx (استيراد كسول)"""
    try:
        from docx import Document
    except ImportError:
        logger.debug("python-docx غير مثبت — تخطي استخراج DOCX")
        return ""
    doc = Document(str(file_path))
    parts: List[str] = [p.text for p in doc.paragraphs if p.text]
    # تضمين النص من الجداول أيضًا
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


def _extract_xlsx_text(file_path: Path) -> str:
    """استخراج نص من XLSX عبر openpyxl (استيراد كسول)"""
    try:
        import openpyxl
    except ImportError:
        logger.debug("openpyxl غير مثبت — تخطي استخراج XLSX")
        return ""
    wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
    rows: List[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows.append(f"### {sheet_name} ###")
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(c) if c is not None else "" for c in row)
            if row_text.strip(" |"):
                rows.append(row_text)
    wb.close()
    return "\n".join(rows[:2000])  # حد أقصى 2000 سطر


def _extract_pptx_text(file_path: Path) -> str:
    """استخراج نص من PPTX عبر python-pptx (استيراد كسول)"""
    try:
        from pptx import Presentation
    except ImportError:
        logger.debug("python-pptx غير مثبت — تخطي استخراج PPTX")
        return ""
    prs = Presentation(str(file_path))
    parts: List[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        parts.append(f"### Slide {slide_num} ###")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
            # تضمين النص من الجداول في الشريحة
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        parts.append(row_text)
    return "\n".join(parts)


# ملاحظة: تم نقل منطق كشف نوع المحتوى إلى metadata_extractor.detect_content_type
# في PR-03. هذه الدالة محفوظة للتوافق مع أي كود خارجي يشير إليها، لكنها
# مجرد وكيل (delegate) للدالة الجديدة.
def _guess_content_type(ext: str) -> str:
    """وكيل للتوافق مع الكود القديم — يُفضّل استخدام detect_content_type مباشرة

    يُرجع نوع المحتوى من قاموس الامتدادات فقط (لا يقرأ magic bytes).
    """
    from .metadata_extractor import _EXT_MIME_FALLBACK
    ext_clean = ext.lower().lstrip(".")
    return _EXT_MIME_FALLBACK.get(ext_clean, "application/octet-stream")
